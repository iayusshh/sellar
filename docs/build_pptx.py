"""
Sellar — Creator Monetization Platform
Beautiful PPTX presentation builder using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0A, 0x0F)   # deep dark
CARD      = RGBColor(0x13, 0x13, 0x1A)   # card bg
BRAND     = RGBColor(0x7C, 0x3A, 0xED)   # purple
BRAND_L   = RGBColor(0xA7, 0x8B, 0xFA)   # light purple
ACCENT    = RGBColor(0x06, 0xB6, 0xD4)   # cyan
ACCENT2   = RGBColor(0xF5, 0x9E, 0x0B)   # amber
SUCCESS   = RGBColor(0x10, 0xB9, 0x81)   # green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD = RGBColor(0x0D, 0x0D, 0x14)

W = Inches(13.33)   # widescreen slide width  (16:9)
H = Inches(7.5)     # widescreen slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ── Helper utilities ────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    fill_bg(s, BG)
    return s

def fill_bg(s, color):
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    sh = s.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    sh.line.width = line_w
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh

def rounded_rect(s, x, y, w, h, fill=None, line=None, line_w=Pt(1.5), radius=0.08):
    from pptx.oxml.ns import qn
    from lxml import etree
    sh = s.shapes.add_shape(5, x, y, w, h)   # 5 = rounded rectangle
    # set corner radius via XML
    sp_elem = sh._element
    prstGeom = sp_elem.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {int(radius * 100000)}')

    sh.line.width = line_w
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh

def grad_rect(s, x, y, w, h, c1, c2, angle=0):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.gradient()
    sh.fill.gradient_angle = angle
    sh.fill.gradient_stops[0].color.rgb = c1
    sh.fill.gradient_stops[1].color.rgb = c2
    sh.line.fill.background()
    return sh

def txb(s, x, y, w, h, text, size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def label_pill(s, x, y, w, text):
    """Small pill label like '✨ Features'"""
    bg_sh = rounded_rect(s, x, y, w, Inches(0.32),
                          fill=RGBColor(0x1A, 0x0D, 0x3A),
                          line=BRAND, line_w=Pt(1))
    txb(s, x + Inches(0.12), y + Inches(0.04), w - Inches(0.24), Inches(0.28),
        text, size=Pt(11), bold=True, color=BRAND_L, align=PP_ALIGN.CENTER)

def divider(s, y, color=RGBColor(0x1E, 0x1E, 0x2E)):
    rect(s, Inches(0.5), y, W - Inches(1), Pt(1), fill=color)

def step_circle(s, x, y, num, c1=BRAND, c2=ACCENT):
    grad_rect(s, x - Inches(0.18), y - Inches(0.02),
              Inches(0.36), Inches(0.36), c1, c2, angle=135)
    txb(s, x - Inches(0.18), y - Inches(0.02),
        Inches(0.36), Inches(0.36), str(num),
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def icon_card(s, x, y, w, h, icon, title, body, icon_bg=None):
    if icon_bg is None:
        icon_bg = RGBColor(0x1A, 0x0D, 0x3A)
    rounded_rect(s, x, y, w, h,
                 fill=CARD,
                 line=RGBColor(0x23, 0x23, 0x33), line_w=Pt(1))
    # icon square
    rounded_rect(s, x + Inches(0.22), y + Inches(0.2),
                 Inches(0.48), Inches(0.48), fill=icon_bg, line=None, line_w=Pt(0))
    txb(s, x + Inches(0.22), y + Inches(0.2),
        Inches(0.48), Inches(0.48), icon,
        size=Pt(18), align=PP_ALIGN.CENTER)
    txb(s, x + Inches(0.22), y + Inches(0.76),
        w - Inches(0.44), Inches(0.32),
        title, size=Pt(12), bold=True, color=WHITE)
    txb(s, x + Inches(0.22), y + Inches(1.1),
        w - Inches(0.44), h - Inches(1.18),
        body, size=Pt(10), color=MUTED, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════════════════
s1 = slide()

# gradient accent stripe top
grad_rect(s1, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

# big radial glow (simulate with a large transparent-edge gradient ellipse)
from pptx.util import Emu
glow = s1.shapes.add_shape(9, Inches(5.5), Inches(-1.5), Inches(8), Inches(6))  # 9=oval
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x14, 0x07, 0x30)
glow.line.fill.background()

# Logo box
logo = s1.shapes.add_shape(5, Inches(0.6), Inches(1.0), Inches(0.7), Inches(0.7))
logo.fill.gradient()
logo.fill.gradient_angle = 135
logo.fill.gradient_stops[0].color.rgb = BRAND
logo.fill.gradient_stops[1].color.rgb = ACCENT
logo.line.fill.background()
txb(s1, Inches(0.6), Inches(1.0), Inches(0.7), Inches(0.7),
    "S", size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# App name
txb(s1, Inches(1.42), Inches(1.06), Inches(2), Inches(0.5),
    "Sellar", size=Pt(26), bold=True, color=WHITE)

# Tagline label
label_pill(s1, Inches(0.6), Inches(1.9), Inches(3.4),
           "🚀  Creator Monetization Platform")

# Main headline
txb(s1, Inches(0.6), Inches(2.4), Inches(6.2), Inches(1.2),
    "Sell Digital Products", size=Pt(46), bold=True, color=WHITE)
txb(s1, Inches(0.6), Inches(3.45), Inches(6.2), Inches(0.7),
    "& Webinars, Effortlessly", size=Pt(40), bold=True, color=BRAND_L)

txb(s1, Inches(0.6), Inches(4.35), Inches(5.8), Inches(0.8),
    "India's creator-first marketplace. Launch your storefront, list digital products,\nschedule webinars, and get paid — all in under 2 minutes.",
    size=Pt(13), color=MUTED, wrap=True)

# KPI pills row
kpis = [("4", "Product Types"), ("₹0", "Setup Cost"), ("INR", "Payments"), ("80%", "Creator Keeps")]
for i, (val, lbl) in enumerate(kpis):
    px = Inches(0.6) + i * Inches(1.7)
    py = Inches(5.35)
    rounded_rect(s1, px, py, Inches(1.55), Inches(0.8),
                 fill=RGBColor(0x10, 0x10, 0x18),
                 line=RGBColor(0x28, 0x18, 0x50), line_w=Pt(1))
    txb(s1, px, py + Inches(0.04), Inches(1.55), Inches(0.38),
        val, size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s1, px, py + Inches(0.42), Inches(1.55), Inches(0.3),
        lbl, size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)

# Right side — storefront preview card
card_x = Inches(7.4)
card_y = Inches(0.9)
card_w = Inches(5.3)
card_h = Inches(6.0)
rounded_rect(s1, card_x, card_y, card_w, card_h,
             fill=DARK_CARD,
             line=RGBColor(0x28, 0x18, 0x50), line_w=Pt(1.5))

# browser bar
rounded_rect(s1, card_x, card_y, card_w, Inches(0.42),
             fill=RGBColor(0x10, 0x10, 0x1A),
             line=RGBColor(0x28, 0x18, 0x50), line_w=Pt(0))
# dots
for di, dc in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
    dot = s1.shapes.add_shape(9, card_x + Inches(0.18) + di * Inches(0.22),
                               card_y + Inches(0.14), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = dc
    dot.line.fill.background()
# URL bar
rounded_rect(s1, card_x + Inches(0.82), card_y + Inches(0.08),
             Inches(3.8), Inches(0.27),
             fill=RGBColor(0x18, 0x18, 0x28), line=None, line_w=Pt(0))
txb(s1, card_x + Inches(0.85), card_y + Inches(0.08),
    Inches(3.7), Inches(0.26),
    "🔒  sellar.app/ayush", size=Pt(9), color=MUTED)

# Creator avatar circle
avatar = s1.shapes.add_shape(9, card_x + Inches(0.3), card_y + Inches(0.62),
                              Inches(0.7), Inches(0.7))
avatar.fill.gradient(); avatar.fill.gradient_angle = 135
avatar.fill.gradient_stops[0].color.rgb = BRAND
avatar.fill.gradient_stops[1].color.rgb = ACCENT
avatar.line.fill.background()
txb(s1, card_x + Inches(0.3), card_y + Inches(0.62), Inches(0.7), Inches(0.7),
    "A", size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s1, card_x + Inches(1.12), card_y + Inches(0.66),
    Inches(3.5), Inches(0.28), "Ayush Anand", size=Pt(12), bold=True, color=WHITE)
txb(s1, card_x + Inches(1.12), card_y + Inches(0.94),
    Inches(3.5), Inches(0.22), "@ayush · Finance & Tech Creator", size=Pt(9), color=MUTED)

# product cards row
products_data = [
    ("📊", "Stock Course", "₹999", BRAND, RGBColor(0x1E,0x1B,0x4B)),
    ("🎯", "IPO Framework", "₹599", ACCENT2, RGBColor(0x06,0x4E,0x3B)),
    ("⚛️", "React Mastery", "₹799", BRAND_L, RGBColor(0x4C,0x1D,0x95)),
]
for i, (icon, title, price, pc, bg) in enumerate(products_data):
    px = card_x + Inches(0.22) + i * Inches(1.64)
    py = card_y + Inches(1.5)
    rounded_rect(s1, px, py, Inches(1.56), Inches(2.3),
                 fill=RGBColor(0x10, 0x10, 0x1A),
                 line=RGBColor(0x22, 0x22, 0x35), line_w=Pt(1))
    # image area
    rounded_rect(s1, px + Inches(0.08), py + Inches(0.1),
                 Inches(1.4), Inches(1.1), fill=bg, line=None, line_w=Pt(0))
    txb(s1, px + Inches(0.08), py + Inches(0.1),
        Inches(1.4), Inches(1.1), icon,
        size=Pt(28), align=PP_ALIGN.CENTER)
    txb(s1, px + Inches(0.1), py + Inches(1.28),
        Inches(1.36), Inches(0.28), title, size=Pt(9), bold=True, color=WHITE)
    txb(s1, px + Inches(0.1), py + Inches(1.58),
        Inches(0.8), Inches(0.26), price, size=Pt(12), bold=True, color=pc)
    # buy button
    grad_rect(s1, px + Inches(0.1), py + Inches(1.92),
              Inches(1.36), Inches(0.26), BRAND, ACCENT, angle=0)
    txb(s1, px + Inches(0.1), py + Inches(1.92),
        Inches(1.36), Inches(0.26), "Buy Now",
        size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# webinar pill
rounded_rect(s1, card_x + Inches(0.22), card_y + Inches(3.95),
             card_w - Inches(0.44), Inches(0.7),
             fill=RGBColor(0x18, 0x11, 0x00),
             line=RGBColor(0x60, 0x3D, 0x00), line_w=Pt(1))
txb(s1, card_x + Inches(0.35), card_y + Inches(4.02),
    Inches(4.5), Inches(0.28),
    "🎙️  LIVE WEBINAR · May 1, 2026", size=Pt(10), bold=True, color=ACCENT2)
txb(s1, card_x + Inches(0.35), card_y + Inches(4.32),
    Inches(4.5), Inches(0.24),
    "Advanced Options Trading  —  ₹799", size=Pt(9), color=MUTED)

# bottom gradient stripe
grad_rect(s1, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLATFORM FEATURES
# ═══════════════════════════════════════════════════════════════════════════
s2 = slide()
grad_rect(s2, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s2, Inches(0.6), Inches(0.3), Inches(2.8), "✨  Platform Features")
txb(s2, Inches(0.6), Inches(0.75), Inches(9), Inches(0.65),
    "Everything You Need to Monetize Your Expertise",
    size=Pt(30), bold=True, color=WHITE)
txb(s2, Inches(0.6), Inches(1.42), Inches(7), Inches(0.4),
    "From storefront creation to payment processing — Sellar handles the entire commerce lifecycle.",
    size=Pt(12), color=MUTED)

features = [
    ("🏪", "Instant Storefront", "Personalized URL at sellar.app/yourhandle in seconds. No design skills needed.", RGBColor(0x1A,0x0D,0x3A)),
    ("💳", "Cashfree Checkout", "Built-in UPI, net banking, cards & wallets via Cashfree Drop-in modal.", RGBColor(0x00,0x1A,0x1F)),
    ("🎙️", "Webinar Scheduling", "Schedule live sessions with capacity limits & auto-generated join tokens.", RGBColor(0x1A,0x10,0x00)),
    ("💰", "Creator Wallet", "Earnings land in wallet instantly. Withdraw via UPI, IMPS, or bank transfer.", RGBColor(0x00,0x1A,0x10)),
    ("📊", "Analytics Dashboard", "Track earnings, transactions, visitors, and demographics in real time.", RGBColor(0x1A,0x0D,0x3A)),
    ("🔐", "Row-Level Security", "PostgreSQL RLS policies — data isolated at DB layer, not client.", RGBColor(0x00,0x10,0x1A)),
    ("📚", "Buyer Library", "Personal library with download links & webinar join buttons after purchase.", RGBColor(0x1A,0x0D,0x3A)),
    ("⭐", "Featured Creators", "Admin-curated Top Creators listing drives organic discovery.", RGBColor(0x1A,0x10,0x00)),
    ("🛡️", "Commission Split", "Atomic PostgreSQL transaction: 80% creator, 20% platform — race-condition safe.", RGBColor(0x00,0x1A,0x10)),
]

cols, rows = 3, 3
cw = Inches(3.9)
ch = Inches(1.65)
cx_start = Inches(0.6)
cy_start = Inches(2.0)
gx = Inches(0.2)
gy = Inches(0.18)

for i, (icon, title, body, ibg) in enumerate(features):
    col = i % cols
    row = i // cols
    cx = cx_start + col * (cw + gx)
    cy = cy_start + row * (ch + gy)
    icon_card(s2, cx, cy, cw, ch, icon, title, body, icon_bg=ibg)

grad_rect(s2, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW TO JOIN
# ═══════════════════════════════════════════════════════════════════════════
s3 = slide()
grad_rect(s3, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s3, Inches(0.6), Inches(0.3), Inches(2.4), "👤  Getting Started")
txb(s3, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "How to Join Sellar", size=Pt(32), bold=True, color=WHITE)
txb(s3, Inches(0.6), Inches(1.38), Inches(8), Inches(0.36),
    "Two account types — Buyer or Creator. The entire process takes under 2 minutes.",
    size=Pt(12), color=MUTED)

# ── Buyer column ────────────────────────────────────────────────────────────
bx = Inches(0.6)
rounded_rect(s3, bx, Inches(1.9), Inches(5.8), Inches(5.3),
             fill=RGBColor(0x00, 0x12, 0x1A),
             line=RGBColor(0x06, 0x60, 0x80), line_w=Pt(1.2))

# column header
grad_rect(s3, bx, Inches(1.9), Inches(5.8), Inches(0.44),
          ACCENT, RGBColor(0x08,0x91,0xB2), angle=0)
txb(s3, bx, Inches(1.9), Inches(5.8), Inches(0.44),
    "🛍️  Buyer Account", size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

buyer_steps = [
    ("1", "Visit /auth/signup", "Enter your email address and choose a secure password."),
    ("2", "Verify Email", "Click the link in your inbox. Supabase Auth handles verification instantly."),
    ("3", "Browse Storefronts", "Discover creators, explore products, purchase with one click."),
    ("4", "Access /library", "All purchased products appear with download links & webinar join buttons."),
]
for i, (num, title, body) in enumerate(buyer_steps):
    sy = Inches(2.5) + i * Inches(1.08)
    # connector line
    if i < len(buyer_steps) - 1:
        rect(s3, bx + Inches(0.58), sy + Inches(0.38),
             Pt(2), Inches(0.78), fill=RGBColor(0x20,0x40,0x55))
    # circle
    circ = s3.shapes.add_shape(9, bx + Inches(0.3), sy + Inches(0.03),
                                Inches(0.36), Inches(0.36))
    circ.fill.gradient(); circ.fill.gradient_angle = 135
    circ.fill.gradient_stops[0].color.rgb = ACCENT
    circ.fill.gradient_stops[1].color.rgb = RGBColor(0x06,0x60,0x80)
    circ.line.fill.background()
    txb(s3, bx + Inches(0.3), sy + Inches(0.03),
        Inches(0.36), Inches(0.36), num,
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s3, bx + Inches(0.82), sy, Inches(4.7), Inches(0.3),
        title, size=Pt(12), bold=True, color=WHITE)
    txb(s3, bx + Inches(0.82), sy + Inches(0.3), Inches(4.7), Inches(0.6),
        body, size=Pt(10), color=MUTED, wrap=True)

# ── Creator column ──────────────────────────────────────────────────────────
cx2 = Inches(6.77)
rounded_rect(s3, cx2, Inches(1.9), Inches(5.9), Inches(5.3),
             fill=RGBColor(0x0E, 0x07, 0x22),
             line=RGBColor(0x4C, 0x1D, 0x95), line_w=Pt(1.2))

grad_rect(s3, cx2, Inches(1.9), Inches(5.9), Inches(0.44),
          BRAND, RGBColor(0x6D,0x28,0xD9), angle=0)
txb(s3, cx2, Inches(1.9), Inches(5.9), Inches(0.44),
    "🚀  Creator Account", size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

creator_steps = [
    ("1", "Visit /become-a-creator", "Available to new users and existing buyers upgrading to creator."),
    ("2", "Choose Your @handle", "Pick a unique handle — becomes your public storefront URL. Real-time availability check."),
    ("3", "Complete Your Profile", "Add display name, bio, avatar, phone, and social links (Instagram, X, LinkedIn, Website)."),
    ("4", "Storefront Goes Live!", "Your public storefront is instantly accessible. Share the link & start listing products."),
]
for i, (num, title, body) in enumerate(creator_steps):
    sy = Inches(2.5) + i * Inches(1.08)
    if i < len(creator_steps) - 1:
        rect(s3, cx2 + Inches(0.58), sy + Inches(0.38),
             Pt(2), Inches(0.78), fill=RGBColor(0x3C,0x18,0x6A))
    circ = s3.shapes.add_shape(9, cx2 + Inches(0.3), sy + Inches(0.03),
                                Inches(0.36), Inches(0.36))
    circ.fill.gradient(); circ.fill.gradient_angle = 135
    circ.fill.gradient_stops[0].color.rgb = BRAND
    circ.fill.gradient_stops[1].color.rgb = ACCENT
    circ.line.fill.background()
    txb(s3, cx2 + Inches(0.3), sy + Inches(0.03),
        Inches(0.36), Inches(0.36), num,
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s3, cx2 + Inches(0.82), sy, Inches(4.8), Inches(0.3),
        title, size=Pt(12), bold=True, color=WHITE)
    txb(s3, cx2 + Inches(0.82), sy + Inches(0.3), Inches(4.8), Inches(0.6),
        body, size=Pt(10), color=MUTED, wrap=True)

grad_rect(s3, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ROLE MATRIX
# ═══════════════════════════════════════════════════════════════════════════
s4 = slide()
grad_rect(s4, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s4, Inches(0.6), Inches(0.3), Inches(2.4), "👥  User Roles")
txb(s4, Inches(0.6), Inches(0.75), Inches(10), Inches(0.6),
    "Three Roles, One Platform", size=Pt(32), bold=True, color=WHITE)
txb(s4, Inches(0.6), Inches(1.38), Inches(9), Inches(0.36),
    "Every user has a clearly scoped role — with route guards enforcing access at the app level.",
    size=Pt(12), color=MUTED)

roles = [
    ("🛍️ Buyer", ACCENT, RGBColor(0x00,0x1A,0x22), RGBColor(0x06,0x60,0x80), [
        ("🔍", "Browse all creator storefronts"),
        ("💳", "Purchase products via Cashfree"),
        ("📚", "Access personal library"),
        ("🎙️", "Join purchased webinars"),
        ("⬆️", "Upgrade to creator anytime"),
        ("🔒", "ProtectedRoute + GeneralRoute guard"),
    ]),
    ("🚀 Creator", BRAND, RGBColor(0x0E,0x07,0x22), RGBColor(0x4C,0x1D,0x95), [
        ("🏪", "Public storefront at /:handle"),
        ("📦", "CRUD products (4 types)"),
        ("💰", "Creator wallet & withdrawals"),
        ("📊", "Earnings analytics dashboard"),
        ("⚙️", "Profile & social link settings"),
        ("🔒", "CreatorRoute guard"),
    ]),
    ("🛡️ Admin / Owner", ACCENT2, RGBColor(0x18,0x10,0x00), RGBColor(0x78,0x35,0x00), [
        ("👥", "Manage all creators"),
        ("💹", "Set per-creator commission rates"),
        ("💸", "Process withdrawal requests"),
        ("📈", "Platform revenue analytics"),
        ("⭐", "Feature & reorder creators"),
        ("🔒", "AdminRoute + OwnerRoute guard"),
    ]),
]

col_w = Inches(3.8)
for i, (role, hdr_c, bg_c, border_c, items) in enumerate(roles):
    rx = Inches(0.6) + i * (col_w + Inches(0.26))
    ry = Inches(1.9)
    rh = Inches(5.3)
    rounded_rect(s4, rx, ry, col_w, rh, fill=bg_c, line=border_c, line_w=Pt(1.5))
    # header
    grad_rect(s4, rx, ry, col_w, Inches(0.52), hdr_c, border_c, angle=0)
    txb(s4, rx, ry, col_w, Inches(0.52),
        role, size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (icon, desc) in enumerate(items):
        iy = ry + Inches(0.66) + j * Inches(0.72)
        rounded_rect(s4, rx + Inches(0.2), iy, col_w - Inches(0.4), Inches(0.58),
                     fill=RGBColor(0x10,0x10,0x18),
                     line=RGBColor(0x22,0x22,0x33), line_w=Pt(0.5))
        txb(s4, rx + Inches(0.32), iy + Inches(0.08),
            Inches(0.3), Inches(0.42), icon, size=Pt(16))
        txb(s4, rx + Inches(0.68), iy + Inches(0.12),
            col_w - Inches(0.88), Inches(0.36),
            desc, size=Pt(10), color=WHITE if j < 5 else MUTED)

grad_rect(s4, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PRODUCT TYPES
# ═══════════════════════════════════════════════════════════════════════════
s5 = slide()
grad_rect(s5, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s5, Inches(0.6), Inches(0.3), Inches(2.4), "📦  Product Catalog")
txb(s5, Inches(0.6), Inches(0.75), Inches(10), Inches(0.6),
    "Four Types of Digital Products", size=Pt(32), bold=True, color=WHITE)
txb(s5, Inches(0.6), Inches(1.38), Inches(9), Inches(0.36),
    "Sell anything from downloadable files to live webinar sessions — Sellar supports every format.",
    size=Pt(12), color=MUTED)

prod_types = [
    ("📥", "Digital Downloads",
     "eBooks · PDFs · Templates · Guides · Spreadsheets\n\nBuyers receive an instant download link upon purchase. Upload any file to Supabase Storage — secure, CDN-served.",
     BRAND, RGBColor(0x1A,0x0D,0x3A)),
    ("🎙️", "Live Webinars",
     "Scheduled · Capacity Limits · Early-Join Window\n\nSet date, time, max seats, duration, and join window. Secure join tokens auto-generated per verified buyer.",
     ACCENT, RGBColor(0x00,0x1A,0x22)),
    ("👤", "1-on-1 Sessions",
     "Consulting · Coaching · Mentoring\n\nPersonal sessions with a booking URL. Buyers get the access link in their library after payment is confirmed.",
     RGBColor(0xEC,0x48,0x99), RGBColor(0x22,0x00,0x12)),
    ("📱", "Telegram Access",
     "Community · Channel · Private Group\n\nShare a private Telegram invite link with verified buyers. Great for premium signal channels and communities.",
     ACCENT2, RGBColor(0x18,0x10,0x00)),
]

tw = Inches(2.9)
th = Inches(4.8)
for i, (icon, title, body, hc, bg) in enumerate(prod_types):
    tx = Inches(0.6) + i * (tw + Inches(0.26))
    ty = Inches(1.95)
    rounded_rect(s5, tx, ty, tw, th, fill=bg,
                 line=hc, line_w=Pt(1.5))
    # icon circle
    circ = s5.shapes.add_shape(9, tx + tw/2 - Inches(0.45), ty + Inches(0.2),
                                Inches(0.9), Inches(0.9))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x10,0x10,0x20)
    circ.line.color.rgb = hc; circ.line.width = Pt(1.5)
    txb(s5, tx + tw/2 - Inches(0.45), ty + Inches(0.2),
        Inches(0.9), Inches(0.9),
        icon, size=Pt(28), align=PP_ALIGN.CENTER)
    txb(s5, tx + Inches(0.2), ty + Inches(1.25),
        tw - Inches(0.4), Inches(0.38),
        title, size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s5, tx + Inches(0.2), ty + Inches(1.72),
        tw - Inches(0.4), th - Inches(1.9),
        body, size=Pt(10), color=MUTED, wrap=True, align=PP_ALIGN.CENTER)

grad_rect(s5, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW TO LIST A PRODUCT
# ═══════════════════════════════════════════════════════════════════════════
s6 = slide()
grad_rect(s6, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s6, Inches(0.6), Inches(0.3), Inches(2.6), "📦  Creator Flow")
txb(s6, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "How to List a Product", size=Pt(32), bold=True, color=WHITE)
txb(s6, Inches(0.6), Inches(1.38), Inches(7.5), Inches(0.36),
    "From dashboard to live in under 2 minutes — CRUD products with image upload, pricing, and scheduling.",
    size=Pt(12), color=MUTED)

list_steps = [
    ("1", "Creator Dashboard", "/creator/products — access from sidebar navigation."),
    ("2", "Click 'Add Product'", "Opens the product creation modal with all required fields."),
    ("3", "Choose Product Type", "Digital Download, Webinar, Session, or Telegram Access."),
    ("4", "Fill Details", "Title, description, price (INR), image upload, content URL."),
    ("5", "Webinars: Schedule", "Date/time, max capacity, duration, early-join & late-join window."),
    ("6", "Publish — It's Live!", "Appears on storefront instantly. Hide or re-activate without deleting."),
]

# Steps (left panel)
for i, (num, title, body) in enumerate(list_steps):
    row = i % 3
    col = i // 3
    sx = Inches(0.6) + col * Inches(3.4)
    sy = Inches(2.0) + row * Inches(1.65)

    rounded_rect(s6, sx, sy, Inches(3.2), Inches(1.5),
                 fill=CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))
    circ = s6.shapes.add_shape(9, sx + Inches(0.15), sy + Inches(0.14),
                                Inches(0.38), Inches(0.38))
    circ.fill.gradient(); circ.fill.gradient_angle = 135
    circ.fill.gradient_stops[0].color.rgb = BRAND
    circ.fill.gradient_stops[1].color.rgb = ACCENT
    circ.line.fill.background()
    txb(s6, sx + Inches(0.15), sy + Inches(0.14),
        Inches(0.38), Inches(0.38), num,
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s6, sx + Inches(0.62), sy + Inches(0.16),
        Inches(2.42), Inches(0.3),
        title, size=Pt(12), bold=True, color=WHITE)
    txb(s6, sx + Inches(0.62), sy + Inches(0.5),
        Inches(2.42), Inches(0.88),
        body, size=Pt(10), color=MUTED, wrap=True)

# Right panel — products list mockup
mx = Inches(7.1)
my = Inches(1.95)
mw = Inches(5.8)
mh = Inches(5.35)
rounded_rect(s6, mx, my, mw, mh,
             fill=DARK_CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1.2))

# browser bar
grad_rect(s6, mx, my, mw, Inches(0.38), RGBColor(0x10,0x10,0x1A), RGBColor(0x10,0x10,0x1A), angle=0)
for di, dc in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
    dot = s6.shapes.add_shape(9, mx + Inches(0.15) + di * Inches(0.2),
                               my + Inches(0.12), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = dc; dot.line.fill.background()
txb(s6, mx + Inches(0.7), my + Inches(0.08),
    Inches(4.5), Inches(0.24),
    "🔒  sellar.app/creator/products", size=Pt(9), color=MUTED)

# Page title
txb(s6, mx + Inches(0.25), my + Inches(0.52), Inches(3.5), Inches(0.3),
    "My Products", size=Pt(14), bold=True, color=WHITE)
txb(s6, mx + Inches(0.25), my + Inches(0.82), Inches(3.5), Inches(0.24),
    "6 products · 3 active", size=Pt(10), color=MUTED)
# add button
grad_rect(s6, mx + Inches(3.8), my + Inches(0.55),
          Inches(1.72), Inches(0.36), BRAND, ACCENT, angle=0)
txb(s6, mx + Inches(3.8), my + Inches(0.55),
    Inches(1.72), Inches(0.36), "+ Add Product",
    size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

prod_rows = [
    ("📊", "Stock Market Course", "Digital · ₹999 · 47 sales", SUCCESS, "Active"),
    ("🎙️", "Options Webinar", "Webinar · ₹799 · Apr 20 · 30/50 seats", ACCENT2, "Upcoming"),
    ("📋", "SaaS Starter Template", "Digital · ₹1,199 · 23 sales", SUCCESS, "Active"),
    ("📱", "Telegram Access", "Telegram · ₹299 · 18 members", ACCENT, "Active"),
]
for i, (icon, name, meta, sc, status) in enumerate(prod_rows):
    ry = my + Inches(1.22) + i * Inches(0.95)
    rounded_rect(s6, mx + Inches(0.2), ry, mw - Inches(0.4), Inches(0.82),
                 fill=RGBColor(0x10,0x10,0x1A),
                 line=RGBColor(0x22,0x22,0x33), line_w=Pt(0.8))
    # icon
    rounded_rect(s6, mx + Inches(0.32), ry + Inches(0.12),
                 Inches(0.58), Inches(0.58),
                 fill=RGBColor(0x1A,0x0D,0x3A), line=None, line_w=Pt(0))
    txb(s6, mx + Inches(0.32), ry + Inches(0.12),
        Inches(0.58), Inches(0.58), icon,
        size=Pt(18), align=PP_ALIGN.CENTER)
    txb(s6, mx + Inches(1.04), ry + Inches(0.1),
        Inches(3.0), Inches(0.3), name, size=Pt(11), bold=True, color=WHITE)
    txb(s6, mx + Inches(1.04), ry + Inches(0.42),
        Inches(3.0), Inches(0.28), meta, size=Pt(9), color=MUTED)
    # status dot + label
    dot = s6.shapes.add_shape(9, mx + mw - Inches(1.5), ry + Inches(0.33),
                               Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = sc; dot.line.fill.background()
    txb(s6, mx + mw - Inches(1.38), ry + Inches(0.26),
        Inches(0.9), Inches(0.24), status, size=Pt(9), bold=True, color=sc)

grad_rect(s6, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HOW TO PURCHASE / PAYMENT FLOW
# ═══════════════════════════════════════════════════════════════════════════
s7 = slide()
grad_rect(s7, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s7, Inches(0.6), Inches(0.3), Inches(2.4), "💳  Buyer Flow")
txb(s7, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "How to Purchase a Product", size=Pt(32), bold=True, color=WHITE)
txb(s7, Inches(0.6), Inches(1.38), Inches(9), Inches(0.36),
    "From storefront to library in under 60 seconds. UPI · Net Banking · Cards · Wallets via Cashfree.",
    size=Pt(12), color=MUTED)

# Payment pipeline arrows
pipeline = [
    ("🏪", "Browse\nStorefront", "Visit /:handle\nPick a product"),
    ("🔑", "Sign In /\nSign Up", "Quick auth if\nnot logged in"),
    ("💳", "Cashfree\nDrop-in", "UPI / Cards /\nNet Banking"),
    ("🔄", "Payment\nVerified", "Webhook + polling\n6 retries · 2.5s"),
    ("📚", "Library\nAccess", "Download links\nand join buttons"),
]

pw = Inches(2.1)
ph = Inches(1.4)
gap = Inches(0.22)
px0 = Inches(0.6)
py0 = Inches(1.92)

for i, (icon, title, desc) in enumerate(pipeline):
    px = px0 + i * (pw + gap + Inches(0.24))
    if i > 0:
        # arrow
        rect(s7, px - Inches(0.28), py0 + ph/2 - Pt(1),
             Inches(0.26), Pt(2), fill=BRAND)
        txb(s7, px - Inches(0.28), py0 + ph/2 - Inches(0.2),
            Inches(0.26), Inches(0.4), "›",
            size=Pt(18), bold=True, color=BRAND_L, align=PP_ALIGN.CENTER)

    is_mid = (i == 2)
    bg_c = CARD if not is_mid else RGBColor(0x14,0x07,0x30)
    bc   = (BRAND if is_mid else RGBColor(0x28,0x18,0x50))
    rounded_rect(s7, px, py0, pw, ph, fill=bg_c, line=bc, line_w=Pt(1.5 if is_mid else 1))
    txb(s7, px, py0 + Inches(0.1), pw, Inches(0.5),
        icon, size=Pt(24), align=PP_ALIGN.CENTER)
    txb(s7, px, py0 + Inches(0.55), pw, Inches(0.46),
        title, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s7, px, py0 + Inches(0.95), pw, Inches(0.38),
        desc, size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)

# Under-the-hood technical row
txb(s7, Inches(0.6), Inches(3.56), Inches(9), Inches(0.3),
    "⚙️  Under the Hood — Payment Pipeline",
    size=Pt(11), bold=True, color=MUTED)

tech_steps = [
    ("🛒 Create Order", "Edge Function", BRAND, RGBColor(0x14,0x07,0x30)),
    ("💳 Cashfree UI", "Drop-in SDK", ACCENT, RGBColor(0x00,0x14,0x1E)),
    ("🔄 Verify Order", "6 retries / 2.5s", ACCENT2, RGBColor(0x18,0x10,0x00)),
    ("💰 Commission Split", "80% Creator + 20% Platform", SUCCESS, RGBColor(0x00,0x14,0x0A)),
    ("📚 Library Update", "React Query invalidation", BRAND_L, RGBColor(0x14,0x07,0x30)),
]
tw2 = Inches(2.3)
th2 = Inches(1.22)
for i, (title, sub, hc, bg) in enumerate(tech_steps):
    tx = Inches(0.6) + i * (tw2 + Inches(0.16))
    ty = Inches(3.96)
    rounded_rect(s7, tx, ty, tw2, th2, fill=bg,
                 line=hc, line_w=Pt(1))
    txb(s7, tx + Inches(0.12), ty + Inches(0.14),
        tw2 - Inches(0.24), Inches(0.38),
        title, size=Pt(11), bold=True, color=WHITE)
    txb(s7, tx + Inches(0.12), ty + Inches(0.56),
        tw2 - Inches(0.24), Inches(0.54),
        sub, size=Pt(9), color=MUTED, wrap=True)

# Checkout mockup
txb(s7, Inches(0.6), Inches(5.36), Inches(4.5), Inches(0.28),
    "Cashfree Drop-In Checkout Modal", size=Pt(11), bold=True, color=MUTED)

cx_mock = Inches(0.6)
cy_mock = Inches(5.7)
cw_mock = Inches(5.8)
ch_mock = Inches(1.58)
rounded_rect(s7, cx_mock, cy_mock, cw_mock, ch_mock,
             fill=DARK_CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))

# order summary
txb(s7, cx_mock + Inches(0.2), cy_mock + Inches(0.12),
    Inches(3.5), Inches(0.26), "Order: IPO Framework Guide  —  ₹599",
    size=Pt(10), bold=True, color=WHITE)

methods = [("📱 UPI", BRAND, True), ("🏦 Net Banking", None, False), ("💳 Card", None, False)]
for i, (label, hc, active) in enumerate(methods):
    mx2 = cx_mock + Inches(0.2) + i * Inches(1.86)
    my2 = cy_mock + Inches(0.48)
    bg2 = RGBColor(0x14,0x07,0x30) if active else RGBColor(0x10,0x10,0x18)
    bc2 = (BRAND if active else RGBColor(0x22,0x22,0x33))
    rounded_rect(s7, mx2, my2, Inches(1.72), Inches(0.36),
                 fill=bg2, line=bc2, line_w=Pt(1))
    txb(s7, mx2, my2, Inches(1.72), Inches(0.36),
        label, size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)

grad_rect(s7, cx_mock + Inches(0.2), cy_mock + Inches(0.96),
          Inches(5.4), Inches(0.42), BRAND, ACCENT, angle=0)
txb(s7, cx_mock + Inches(0.2), cy_mock + Inches(0.96),
    Inches(5.4), Inches(0.42), "Pay ₹599 Securely  🔐",
    size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Library mockup
txb(s7, Inches(6.65), Inches(5.36), Inches(4.5), Inches(0.28),
    "Buyer Library (/library)", size=Pt(11), bold=True, color=MUTED)

lx = Inches(6.65)
ly = Inches(5.7)
lw = Inches(6.45)
lh = Inches(1.58)
rounded_rect(s7, lx, ly, lw, lh,
             fill=DARK_CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))

lib_items = [
    ("📈", "IPO Framework", "⬇ Download", SUCCESS),
    ("📊", "Stock Course", "⬇ Download", SUCCESS),
    ("🎙️", "Options Webinar", "📅 Join Now", ACCENT2),
    ("📱", "Telegram", "🔗 Join Group", ACCENT),
]
for i, (icon, name, action, ac) in enumerate(lib_items):
    ix = lx + Inches(0.2) + i * Inches(1.56)
    iy = ly + Inches(0.1)
    rounded_rect(s7, ix, iy, Inches(1.44), lh - Inches(0.2),
                 fill=RGBColor(0x10,0x10,0x1A),
                 line=RGBColor(0x22,0x22,0x33), line_w=Pt(0.8))
    txb(s7, ix, iy + Inches(0.08), Inches(1.44), Inches(0.44),
        icon, size=Pt(20), align=PP_ALIGN.CENTER)
    txb(s7, ix + Inches(0.06), iy + Inches(0.52),
        Inches(1.32), Inches(0.3), name, size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rounded_rect(s7, ix + Inches(0.12), iy + Inches(0.86),
                 Inches(1.2), Inches(0.28),
                 fill=RGBColor(0x05,0x14,0x0A) if ac == SUCCESS else
                      (RGBColor(0x18,0x10,0x00) if ac == ACCENT2 else RGBColor(0x00,0x14,0x22)),
                 line=ac, line_w=Pt(1))
    txb(s7, ix + Inches(0.12), iy + Inches(0.86),
        Inches(1.2), Inches(0.28), action,
        size=Pt(8), bold=True, color=ac, align=PP_ALIGN.CENTER)

grad_rect(s7, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CREATOR DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════
s8 = slide()
grad_rect(s8, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s8, Inches(0.6), Inches(0.3), Inches(2.8), "📊  Creator Dashboard")
txb(s8, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "Creator Dashboard & Wallet", size=Pt(32), bold=True, color=WHITE)
txb(s8, Inches(0.6), Inches(1.38), Inches(9), Inches(0.36),
    "Track earnings, manage products, and request withdrawals — all from one unified dashboard.",
    size=Pt(12), color=MUTED)

# KPI cards
kpis_dash = [
    ("₹28,450", "Available Balance", SUCCESS, "↑ Ready to withdraw"),
    ("₹12,800", "Monthly Earnings", WHITE, "↑ +23% vs last month"),
    ("₹1,24,600", "Lifetime Earnings", WHITE, "Total since joining"),
    ("₹5,000", "Pending Withdrawals", ACCENT2, "Processing"),
]
kw = Inches(2.88)
kh = Inches(1.08)
for i, (val, lbl, vc, sub) in enumerate(kpis_dash):
    kx = Inches(0.6) + i * (kw + Inches(0.2))
    ky = Inches(1.9)
    rounded_rect(s8, kx, ky, kw, kh,
                 fill=CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))
    txb(s8, kx + Inches(0.18), ky + Inches(0.08),
        kw - Inches(0.36), Inches(0.26),
        lbl, size=Pt(10), color=MUTED)
    txb(s8, kx + Inches(0.18), ky + Inches(0.34),
        kw - Inches(0.36), Inches(0.44),
        val, size=Pt(20), bold=True, color=vc)
    txb(s8, kx + Inches(0.18), ky + Inches(0.78),
        kw - Inches(0.36), Inches(0.22),
        sub, size=Pt(9), color=SUCCESS if "↑" in sub else MUTED)

# Earnings chart (bar chart visual)
chart_x = Inches(0.6)
chart_y = Inches(3.15)
chart_w = Inches(7.0)
chart_h = Inches(1.9)
rounded_rect(s8, chart_x, chart_y, chart_w, chart_h,
             fill=CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))
txb(s8, chart_x + Inches(0.2), chart_y + Inches(0.1),
    Inches(4), Inches(0.28), "Earnings — Last 7 Months",
    size=Pt(11), bold=True, color=MUTED)

months = ["Oct", "Nov", "Dec", "Jan", "Feb", "Mar", "Apr"]
heights = [0.40, 0.55, 0.46, 0.70, 0.62, 0.85, 1.0]
bar_w = Inches(0.58)
bar_max_h = Inches(1.18)
bar_x0 = chart_x + Inches(0.5)
bar_base = chart_y + chart_h - Inches(0.32)

for i, (m, h) in enumerate(zip(months, heights)):
    bh = bar_max_h * h
    bx = bar_x0 + i * Inches(0.86)
    by = bar_base - bh
    is_last = (i == 6)
    c1 = BRAND if not is_last else BRAND_L
    c2 = ACCENT if not is_last else ACCENT
    grad_rect(s8, bx, by, bar_w, bh, c1, c2, angle=90)
    txb(s8, bx, bar_base + Inches(0.04), bar_w, Inches(0.22),
        m, size=Pt(9), color=(BRAND_L if is_last else MUTED), align=PP_ALIGN.CENTER)

# Recent transactions table
tx_x = Inches(0.6)
tx_y = Inches(5.22)
tx_w = Inches(7.0)
tx_h = Inches(1.98)
rounded_rect(s8, tx_x, tx_y, tx_w, tx_h,
             fill=CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))
txb(s8, tx_x + Inches(0.2), tx_y + Inches(0.08),
    Inches(3), Inches(0.28), "Recent Transactions",
    size=Pt(11), bold=True, color=MUTED)

# table header
for hdr, cx3 in [("Product", Inches(0.2)), ("Buyer", Inches(2.2)),
                  ("Amount", Inches(3.8)), ("Status", Inches(4.8)), ("Date", Inches(5.8))]:
    txb(s8, tx_x + cx3, tx_y + Inches(0.38), Inches(0.9), Inches(0.24),
        hdr, size=Pt(9), color=MUTED, bold=True)

rows_data = [
    ("Stock Market Course", "Rahul M.", "₹799", "Completed", "Apr 14", SUCCESS),
    ("SaaS Starter", "Sneha K.", "₹959", "Completed", "Apr 13", SUCCESS),
    ("Options Webinar", "Vikram S.", "₹639", "Pending", "Apr 12", ACCENT2),
]
for ri, (prod, buyer, amt, stat, date, sc) in enumerate(rows_data):
    ry2 = tx_y + Inches(0.66) + ri * Inches(0.38)
    if ri % 2 == 0:
        rect(s8, tx_x + Inches(0.1), ry2 - Inches(0.03),
             tx_w - Inches(0.2), Inches(0.36),
             fill=RGBColor(0x10,0x10,0x18))
    for val, cx3 in [(prod, Inches(0.2)), (buyer, Inches(2.2))]:
        txb(s8, tx_x + cx3, ry2, Inches(1.8), Inches(0.3), val, size=Pt(10), color=WHITE)
    txb(s8, tx_x + Inches(3.8), ry2, Inches(0.9), Inches(0.3),
        amt, size=Pt(10), bold=True, color=SUCCESS)
    rounded_rect(s8, tx_x + Inches(4.7), ry2 + Inches(0.02),
                 Inches(0.88), Inches(0.26),
                 fill=RGBColor(0x00,0x14,0x0A) if sc == SUCCESS else RGBColor(0x18,0x10,0x00),
                 line=sc, line_w=Pt(0.8))
    txb(s8, tx_x + Inches(4.7), ry2 + Inches(0.02),
        Inches(0.88), Inches(0.26), stat,
        size=Pt(9), bold=True, color=sc, align=PP_ALIGN.CENTER)
    txb(s8, tx_x + Inches(5.7), ry2, Inches(0.9), Inches(0.3),
        date, size=Pt(9), color=MUTED)

# Wallet card (right)
wl_x = Inches(7.88)
wl_y = Inches(1.9)
wl_w = Inches(5.0)
wl_h = Inches(2.48)
rounded_rect(s8, wl_x, wl_y, wl_w, wl_h,
             fill=RGBColor(0x10,0x07,0x28),
             line=BRAND, line_w=Pt(1.5))
txb(s8, wl_x + Inches(0.22), wl_y + Inches(0.12),
    Inches(4.5), Inches(0.24), "Available Balance",
    size=Pt(10), color=MUTED)
txb(s8, wl_x + Inches(0.22), wl_y + Inches(0.38),
    Inches(4.5), Inches(0.6), "₹28,450",
    size=Pt(30), bold=True, color=BRAND_L)
txb(s8, wl_x + Inches(0.22), wl_y + Inches(1.0),
    Inches(4.5), Inches(0.24), "Ready to withdraw · INR",
    size=Pt(10), color=MUTED)

wl_stats = [("Monthly", "₹12,800"), ("Lifetime", "₹1,24,600"), ("Commission", "20%")]
for i, (lbl, val) in enumerate(wl_stats):
    wx2 = wl_x + Inches(0.22) + i * Inches(1.58)
    txb(s8, wx2, wl_y + Inches(1.38), Inches(1.4), Inches(0.24),
        lbl, size=Pt(9), color=MUTED)
    txb(s8, wx2, wl_y + Inches(1.62), Inches(1.4), Inches(0.36),
        val, size=Pt(14), bold=True, color=WHITE)

# Withdrawal methods
txb(s8, wl_x + Inches(0.22), wl_y + Inches(2.08),
    Inches(4.5), Inches(0.24), "Withdraw via: UPI · IMPS/NEFT · Bank Transfer",
    size=Pt(10), color=ACCENT)

# Sidebar nav
sb_x = Inches(7.88)
sb_y = Inches(4.55)
sb_w = Inches(5.0)
sb_h = Inches(2.6)
rounded_rect(s8, sb_x, sb_y, sb_w, sb_h,
             fill=CARD, line=RGBColor(0x22,0x22,0x33), line_w=Pt(1))
txb(s8, sb_x + Inches(0.2), sb_y + Inches(0.1),
    Inches(4.5), Inches(0.26), "Creator Navigation",
    size=Pt(11), bold=True, color=MUTED)

nav_items = [
    ("📊 Dashboard", True), ("📦 My Products", False),
    ("💰 Wallet", False), ("⚙️ Settings", False),
    ("🏪 View Storefront", False),
]
for i, (label, active) in enumerate(nav_items):
    ny = sb_y + Inches(0.42) + i * Inches(0.38)
    if active:
        rect(s8, sb_x + Inches(0.12), ny, sb_w - Inches(0.24), Inches(0.32),
             fill=RGBColor(0x1A,0x0D,0x3A))
    txb(s8, sb_x + Inches(0.22), ny + Inches(0.04),
        Inches(4.2), Inches(0.26),
        label, size=Pt(11), color=(BRAND_L if active else MUTED), bold=active)

grad_rect(s8, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ADMIN PORTAL
# ═══════════════════════════════════════════════════════════════════════════
s9 = slide()
grad_rect(s9, 0, 0, W, Inches(0.06), ACCENT2, RGBColor(0xD9,0x77,0x06), angle=0)

label_pill(s9, Inches(0.6), Inches(0.3), Inches(2.4), "🛡️  Admin Portal")
txb(s9, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "Full Platform Visibility & Control", size=Pt(32), bold=True, color=WHITE)
txb(s9, Inches(0.6), Inches(1.38), Inches(9), Inches(0.36),
    "8-tab admin portal with creator management, commission control, withdrawal processing, and analytics.",
    size=Pt(12), color=MUTED)

# 8 admin tabs
tabs = ["📊 Overview", "👥 Creators", "📦 Products", "👤 Users",
        "💸 Withdrawals", "💹 Commission", "📈 Traffic", "🗺️ Demographics"]
tab_w = Inches(1.48)
for i, tab in enumerate(tabs):
    tx = Inches(0.6) + i * (tab_w + Inches(0.08))
    active = (i == 0)
    rounded_rect(s9, tx, Inches(1.9), tab_w, Inches(0.38),
                 fill=(RGBColor(0x1A,0x10,0x00) if active else CARD),
                 line=(ACCENT2 if active else RGBColor(0x28,0x28,0x38)),
                 line_w=Pt(1.2 if active else 0.8))
    txb(s9, tx, Inches(1.9), tab_w, Inches(0.38),
        tab, size=Pt(8.5), bold=active,
        color=(ACCENT2 if active else MUTED), align=PP_ALIGN.CENTER)

# Platform KPIs
admin_kpis = [
    ("₹2,48,000", "Platform Revenue", ACCENT2, "↑ +18% this month"),
    ("124", "Active Creators", WHITE, "12 featured"),
    ("1,847", "Total Purchases", WHITE, "↑ +34% this month"),
    ("₹48,200", "Pending Withdrawals", RGBColor(0xEF,0x44,0x44), "8 requests"),
]
kw3 = Inches(2.88)
kh3 = Inches(1.02)
for i, (val, lbl, vc, sub) in enumerate(admin_kpis):
    kx = Inches(0.6) + i * (kw3 + Inches(0.2))
    ky = Inches(2.42)
    rounded_rect(s9, kx, ky, kw3, kh3,
                 fill=CARD, line=RGBColor(0x28,0x28,0x38), line_w=Pt(1))
    txb(s9, kx + Inches(0.15), ky + Inches(0.06),
        kw3 - Inches(0.3), Inches(0.24), lbl, size=Pt(9), color=MUTED)
    txb(s9, kx + Inches(0.15), ky + Inches(0.3),
        kw3 - Inches(0.3), Inches(0.4), val, size=Pt(20), bold=True, color=vc)
    txb(s9, kx + Inches(0.15), ky + Inches(0.72),
        kw3 - Inches(0.3), Inches(0.22), sub, size=Pt(9), color=MUTED)

# Creator management table
table_x = Inches(0.6)
table_y = Inches(3.58)
table_w = Inches(12.1)
table_h = Inches(2.0)
rounded_rect(s9, table_x, table_y, table_w, table_h,
             fill=CARD, line=RGBColor(0x28,0x28,0x38), line_w=Pt(1))
txb(s9, table_x + Inches(0.2), table_y + Inches(0.1),
    Inches(4), Inches(0.26), "Creator Management",
    size=Pt(11), bold=True, color=MUTED)

# headers
hdrs = [("Creator", Inches(0.2)), ("Handle", Inches(2.2)), ("Products", Inches(3.8)),
        ("Commission", Inches(5.0)), ("Featured", Inches(6.3)), ("Actions", Inches(7.5))]
for hdr, hx in hdrs:
    txb(s9, table_x + hx, table_y + Inches(0.4), Inches(1.2), Inches(0.24),
        hdr, size=Pt(9), color=MUTED, bold=True)

creator_rows = [
    ("A", "Ayush Anand", "@ayush", "6 products", "20%", "⭐ Yes",
     RGBColor(0x7C,0x3A,0xED), RGBColor(0x06,0xB6,0xD4)),
    ("P", "Priya Sharma", "@priya", "4 products", "15%", "⭐ Yes",
     RGBColor(0xEC,0x48,0x99), RGBColor(0xF5,0x9E,0x0B)),
    ("R", "Ravi Mehta", "@ravi", "3 products", "18%", "—",
     RGBColor(0x10,0xB9,0x81), RGBColor(0x06,0xB6,0xD4)),
]
for ri, (init, name, handle, prods, comm, feat, c1, c2) in enumerate(creator_rows):
    ry3 = table_y + Inches(0.68) + ri * Inches(0.38)
    if ri % 2 == 0:
        rect(s9, table_x + Inches(0.1), ry3 - Inches(0.02),
             table_w - Inches(0.2), Inches(0.34),
             fill=RGBColor(0x10,0x10,0x18))
    # avatar
    av = s9.shapes.add_shape(9, table_x + Inches(0.2), ry3 + Inches(0.01),
                              Inches(0.28), Inches(0.28))
    av.fill.gradient(); av.fill.gradient_angle = 135
    av.fill.gradient_stops[0].color.rgb = c1
    av.fill.gradient_stops[1].color.rgb = c2
    av.line.fill.background()
    txb(s9, table_x + Inches(0.2), ry3 + Inches(0.01),
        Inches(0.28), Inches(0.28), init,
        size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s9, table_x + Inches(0.56), ry3 + Inches(0.04),
        Inches(1.56), Inches(0.26), name, size=Pt(10), color=WHITE)
    txb(s9, table_x + Inches(2.2), ry3 + Inches(0.04),
        Inches(1.5), Inches(0.26), handle, size=Pt(10), color=BRAND_L)
    txb(s9, table_x + Inches(3.8), ry3 + Inches(0.04),
        Inches(1.1), Inches(0.26), prods, size=Pt(10), color=MUTED)
    # commission badge
    rounded_rect(s9, table_x + Inches(5.0), ry3 + Inches(0.04),
                 Inches(0.7), Inches(0.24),
                 fill=RGBColor(0x1A,0x0D,0x3A), line=BRAND_L, line_w=Pt(0.8))
    txb(s9, table_x + Inches(5.0), ry3 + Inches(0.04),
        Inches(0.7), Inches(0.24), comm,
        size=Pt(9), bold=True, color=BRAND_L, align=PP_ALIGN.CENTER)
    feat_c = SUCCESS if feat != "—" else MUTED
    txb(s9, table_x + Inches(6.3), ry3 + Inches(0.04),
        Inches(1.1), Inches(0.26), feat, size=Pt(10), color=feat_c)
    # action buttons
    rounded_rect(s9, table_x + Inches(7.5), ry3 + Inches(0.04),
                 Inches(0.6), Inches(0.24),
                 fill=RGBColor(0x10,0x10,0x18), line=RGBColor(0x30,0x30,0x45), line_w=Pt(0.8))
    txb(s9, table_x + Inches(7.5), ry3 + Inches(0.04),
        Inches(0.6), Inches(0.24), "Edit",
        size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)
    rounded_rect(s9, table_x + Inches(8.2), ry3 + Inches(0.04),
                 Inches(0.7), Inches(0.24),
                 fill=RGBColor(0x1A,0x04,0x04), line=RGBColor(0xEF,0x44,0x44), line_w=Pt(0.8))
    txb(s9, table_x + Inches(8.2), ry3 + Inches(0.04),
        Inches(0.7), Inches(0.24), "Remove",
        size=Pt(9), color=RGBColor(0xEF,0x44,0x44), align=PP_ALIGN.CENTER)

# Withdrawal request preview
wr_x = Inches(0.6)
wr_y = Inches(5.72)
rounded_rect(s9, wr_x, wr_y, Inches(12.1), Inches(1.52),
             fill=CARD, line=RGBColor(0x28,0x28,0x38), line_w=Pt(1))
txb(s9, wr_x + Inches(0.2), wr_y + Inches(0.1),
    Inches(4), Inches(0.26), "Pending Withdrawal Requests",
    size=Pt(11), bold=True, color=MUTED)

wr_items = [
    ("Ayush Anand", "UPI · Apr 13", "₹15,000", ACCENT2),
    ("Priya Sharma", "Bank · Apr 12", "₹8,500", ACCENT2),
    ("Ravi Mehta", "IMPS · Apr 11", "₹4,200", ACCENT2),
]
for i, (name, meta, amt, _) in enumerate(wr_items):
    wx3 = wr_x + Inches(0.2) + i * Inches(4.0)
    wy3 = wr_y + Inches(0.44)
    rounded_rect(s9, wx3, wy3, Inches(3.8), Inches(0.82),
                 fill=RGBColor(0x18,0x10,0x00),
                 line=RGBColor(0x60,0x3D,0x00), line_w=Pt(1))
    txb(s9, wx3 + Inches(0.15), wy3 + Inches(0.08),
        Inches(2.0), Inches(0.28), name, size=Pt(11), bold=True, color=WHITE)
    txb(s9, wx3 + Inches(0.15), wy3 + Inches(0.38),
        Inches(2.0), Inches(0.24), meta, size=Pt(9), color=MUTED)
    txb(s9, wx3 + Inches(2.2), wy3 + Inches(0.18),
        Inches(0.9), Inches(0.36), amt, size=Pt(14), bold=True, color=WHITE)
    grad_rect(s9, wx3 + Inches(3.12), wy3 + Inches(0.18),
              Inches(0.5), Inches(0.28), SUCCESS, RGBColor(0x05,0x80,0x40), angle=0)
    txb(s9, wx3 + Inches(3.12), wy3 + Inches(0.18),
        Inches(0.5), Inches(0.28), "✓",
        size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

grad_rect(s9, 0, H - Inches(0.06), W, Inches(0.06), ACCENT2, RGBColor(0xD9,0x77,0x06), angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECURITY & ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s10 = slide()
grad_rect(s10, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s10, Inches(0.6), Inches(0.3), Inches(3.0), "🔐  Security & Architecture")
txb(s10, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "Built Secure by Default", size=Pt(32), bold=True, color=WHITE)
txb(s10, Inches(0.6), Inches(1.38), Inches(9), Inches(0.36),
    "Security enforced at the database layer — not just the client. Zero trust by design.",
    size=Pt(12), color=MUTED)

security = [
    ("🛡️", "Row-Level Security", "PostgreSQL RLS on every table. Creators only access their own data — enforced by the DB, not client code.", RGBColor(0x1A,0x0D,0x3A)),
    ("⚡", "Atomic Commission Split", "Single PostgreSQL stored procedure: commission deduction + creator wallet credit in one race-condition-safe transaction.", RGBColor(0x00,0x1A,0x10)),
    ("🔑", "JWT Authentication", "Supabase Auth issues signed JWTs. Edge Functions verify tokens server-side before any payment operation.", RGBColor(0x1A,0x0D,0x3A)),
    ("🔗", "Secure Join Tokens", "Webinar URLs generated server-side. Only verified purchasers within the join window receive valid tokens.", RGBColor(0x00,0x1A,0x10)),
    ("📡", "Webhook Verification", "Cashfree webhook signature validation prevents spoofed payment notifications.", RGBColor(0x1A,0x0D,0x3A)),
    ("🏦", "RBI-Compliant Payments", "Cashfree is RBI-licensed. UPI, net banking, and card payments are fully regulated and compliant.", RGBColor(0x00,0x1A,0x10)),
]

cols_sec = 3
cw_sec = Inches(3.9)
ch_sec = Inches(1.52)
cx0 = Inches(0.6)
cy0 = Inches(1.92)
gx_sec = Inches(0.2)
gy_sec = Inches(0.16)
for i, (icon, title, body, ibg) in enumerate(security):
    col = i % cols_sec
    row = i // cols_sec
    cx4 = cx0 + col * (cw_sec + gx_sec)
    cy4 = cy0 + row * (ch_sec + gy_sec)
    rounded_rect(s10, cx4, cy4, cw_sec, ch_sec,
                 fill=ibg, line=BRAND, line_w=Pt(1))
    txb(s10, cx4 + Inches(0.18), cy4 + Inches(0.1),
        Inches(0.4), Inches(0.44), icon, size=Pt(22))
    txb(s10, cx4 + Inches(0.66), cy4 + Inches(0.12),
        cw_sec - Inches(0.82), Inches(0.3),
        title, size=Pt(12), bold=True, color=WHITE)
    txb(s10, cx4 + Inches(0.18), cy4 + Inches(0.52),
        cw_sec - Inches(0.36), ch_sec - Inches(0.62),
        body, size=Pt(10), color=MUTED, wrap=True)

# Architecture diagram
arch_y = Inches(5.2)
txb(s10, Inches(0.6), arch_y, Inches(9), Inches(0.28),
    "System Architecture", size=Pt(11), bold=True, color=MUTED)

arch_nodes = [
    ("🌐 Browser SPA", "React 18 · Vite · Tailwind", BRAND, RGBColor(0x10,0x07,0x28)),
    ("▲ Vercel CDN", "Global edge · Static deploy", RGBColor(0xFF,0xFF,0xFF), RGBColor(0x10,0x10,0x18)),
    ("🐘 Supabase", "Auth · PostgreSQL · Storage · Edge Fns", SUCCESS, RGBColor(0x00,0x14,0x0A)),
    ("💳 Cashfree", "UPI · Cards · Net Banking · RBI", ACCENT, RGBColor(0x00,0x10,0x1A)),
]
aw = Inches(2.7)
ah = Inches(1.5)
ax0 = Inches(0.6)
ay0 = arch_y + Inches(0.38)
for i, (title, sub, bc, bg) in enumerate(arch_nodes):
    ax = ax0 + i * (aw + Inches(0.4))
    rounded_rect(s10, ax, ay0, aw, ah, fill=bg, line=bc, line_w=Pt(1.5))
    txb(s10, ax + Inches(0.15), ay0 + Inches(0.18),
        aw - Inches(0.3), Inches(0.42),
        title, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s10, ax + Inches(0.15), ay0 + Inches(0.64),
        aw - Inches(0.3), Inches(0.56),
        sub, size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)
    if i < len(arch_nodes) - 1:
        # arrow
        arrow_x = ax + aw + Inches(0.08)
        txb(s10, arrow_x, ay0 + Inches(0.55),
            Inches(0.28), Inches(0.38),
            "⇄", size=Pt(18), bold=True, color=BRAND_L, align=PP_ALIGN.CENTER)

grad_rect(s10, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════
s11 = slide()
grad_rect(s11, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

label_pill(s11, Inches(0.6), Inches(0.3), Inches(2.6), "⚙️  Technology Stack")
txb(s11, Inches(0.6), Inches(0.75), Inches(9), Inches(0.6),
    "Modern Cloud-Native Stack", size=Pt(32), bold=True, color=WHITE)
txb(s11, Inches(0.6), Inches(1.38), Inches(9.5), Inches(0.36),
    "No custom backend — Supabase handles auth, database, storage, and serverless functions. Deployed on Vercel.",
    size=Pt(12), color=MUTED)

tech = [
    ("⚛️", "React 18", "Frontend SPA", BRAND, RGBColor(0x1A,0x0D,0x3A)),
    ("TS", "TypeScript", "Type Safety", RGBColor(0x31,0x78,0xC6), RGBColor(0x00,0x10,0x28)),
    ("⚡", "Vite + SWC", "Build Tool", RGBColor(0xF5,0x9E,0x0B), RGBColor(0x18,0x10,0x00)),
    ("TW", "Tailwind CSS", "Styling", ACCENT, RGBColor(0x00,0x1A,0x22)),
    ("🎯", "Radix UI", "Components", BRAND_L, RGBColor(0x1A,0x0D,0x3A)),
    ("🐘", "Supabase", "BaaS Platform", SUCCESS, RGBColor(0x00,0x14,0x0A)),
    ("🔒", "PostgreSQL 15", "DB + RLS", RGBColor(0x33,0x6B,0x91), RGBColor(0x00,0x0E,0x18)),
    ("🦕", "Deno", "Edge Functions", RGBColor(0xFF,0xFF,0xFF), RGBColor(0x10,0x10,0x18)),
    ("💳", "Cashfree", "Payments", RGBColor(0x0A,0xB0,0x8A), RGBColor(0x00,0x14,0x10)),
    ("▲", "Vercel", "Deployment", RGBColor(0xFF,0xFF,0xFF), RGBColor(0x10,0x10,0x18)),
]

tw3 = Inches(2.3)
th3 = Inches(1.75)
cols_t = 5
gx_t = Inches(0.24)
gy_t = Inches(0.22)
tx0 = Inches(0.7)
ty0 = Inches(1.9)
for i, (logo, name, role, bc, bg) in enumerate(tech):
    col = i % cols_t
    row = i // cols_t
    tx2 = tx0 + col * (tw3 + gx_t)
    ty2 = ty0 + row * (th3 + gy_t)
    rounded_rect(s11, tx2, ty2, tw3, th3, fill=bg, line=bc, line_w=Pt(1.5))
    txb(s11, tx2, ty2 + Inches(0.18), tw3, Inches(0.56),
        logo, size=Pt(26), bold=True, color=bc, align=PP_ALIGN.CENTER)
    txb(s11, tx2, ty2 + Inches(0.8), tw3, Inches(0.38),
        name, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s11, tx2, ty2 + Inches(1.18), tw3, Inches(0.3),
        role, size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

# DB Schema note
schema_x = Inches(0.6)
schema_y = Inches(6.38)
rounded_rect(s11, schema_x, schema_y, Inches(12.1), Inches(0.88),
             fill=CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))
txb(s11, schema_x + Inches(0.2), schema_y + Inches(0.1),
    Inches(11.5), Inches(0.28), "Database Schema Tables:",
    size=Pt(10), bold=True, color=MUTED)
txb(s11, schema_x + Inches(0.2), schema_y + Inches(0.42),
    Inches(11.5), Inches(0.28),
    "users  ·  wallets  ·  transactions  ·  products  ·  purchases  ·  clients  ·  visits  ·  featured_creators",
    size=Pt(10), color=BRAND_L)

grad_rect(s11, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CTA / CLOSE
# ═══════════════════════════════════════════════════════════════════════════
s12 = slide()
grad_rect(s12, 0, 0, W, Inches(0.06), BRAND, ACCENT, angle=0)

# big glow
glow2 = s12.shapes.add_shape(9, Inches(3), Inches(1), Inches(8), Inches(6))
glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(0x14,0x07,0x30)
glow2.line.fill.background()

label_pill(s12, W/2 - Inches(1.5), Inches(1.0), Inches(3), "🚀  Ready to Start?")

txb(s12, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.0),
    "Launch Your Creator", size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s12, Inches(1.5), Inches(2.52), Inches(10.3), Inches(0.8),
    "Storefront Today", size=Pt(48), bold=True, color=BRAND_L, align=PP_ALIGN.CENTER)

txb(s12, Inches(2.5), Inches(3.45), Inches(8.3), Inches(0.5),
    "Join hundreds of Indian creators selling digital products and webinars.\nYour storefront is ready in under 2 minutes — no technical skills required.",
    size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

# CTA buttons
grad_rect(s12, Inches(4.1), Inches(4.22), Inches(2.4), Inches(0.58),
          BRAND, ACCENT, angle=0)
txb(s12, Inches(4.1), Inches(4.22), Inches(2.4), Inches(0.58),
    "Become a Creator →", size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rounded_rect(s12, Inches(6.78), Inches(4.22), Inches(2.4), Inches(0.58),
             fill=RGBColor(0x10,0x10,0x18),
             line=RGBColor(0x28,0x18,0x50), line_w=Pt(1.5))
txb(s12, Inches(6.78), Inches(4.22), Inches(2.4), Inches(0.58),
    "Browse Products", size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Stats row
stats = [("₹0", "Setup cost"), ("80%", "Creator keeps"), ("2 min", "To launch"), ("4", "Product types"), ("INR", "Local payments")]
stat_w = Inches(2.2)
stat_x0 = W/2 - (len(stats) * stat_w + (len(stats)-1) * Inches(0.1)) / 2
for i, (val, lbl) in enumerate(stats):
    sx = stat_x0 + i * (stat_w + Inches(0.1))
    sy = Inches(5.08)
    rounded_rect(s12, sx, sy, stat_w, Inches(1.1),
                 fill=CARD, line=RGBColor(0x28,0x18,0x50), line_w=Pt(1))
    txb(s12, sx, sy + Inches(0.06), stat_w, Inches(0.52),
        val, size=Pt(26), bold=True, color=BRAND_L, align=PP_ALIGN.CENTER)
    txb(s12, sx, sy + Inches(0.62), stat_w, Inches(0.34),
        lbl, size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

# Footer
txb(s12, Inches(0.6), Inches(6.72), Inches(12), Inches(0.5),
    "Sellar — Creator Monetization Platform  ·  Built with React + Supabase + Cashfree  ·  Made in India 🇮🇳  ·  © 2026 Ayush Anand",
    size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

grad_rect(s12, 0, H - Inches(0.06), W, Inches(0.06), BRAND, ACCENT, angle=0)


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
out = "/Users/iayusshh/Work/Sellar/docs/Sellar_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
